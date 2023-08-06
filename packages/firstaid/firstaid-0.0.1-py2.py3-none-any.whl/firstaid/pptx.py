from io import BytesIO

import pandas as pd
from PIL import Image
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Cm, Pt


class PPTX():
    def __init__(self):
        self.prs = Presentation()

        # slide dimensions
        self.SLIDE_HEIGHT = self.prs.slide_height.pt
        self.SLIDE_WIDTH = self.prs.slide_width.pt
        self.Y = Cm(4.45).pt
        self.X = Cm(1.27).pt
        self.TWO_ITEM_X = Cm(12.91).pt
        self.TWO_ITEM_WIDTH = Cm(11.22).pt
        self.ONE_ITEM_WIDTH = Cm(22.86).pt
        self.ITEM_HEIGHT = Cm(12.57).pt
        self.TABLE_ROW_HEIGHT = Cm(1)

    def add_bulletpoint(self, item, shapes, placeholder, layout_number=None):
        """
        add strings to slide
        :param item: list of strings. ex) ['first bullet', 'second bullet', 'third bullet']
        """
        body_shape = shapes.placeholders[placeholder]
        tf = body_shape.text_frame
        for i, row in enumerate(item):
            if i == 0:
                tf.text = row
            else:
                p = tf.add_paragraph()
                p.text = row
        return shapes

    def object_size(self, placeholder, layout_number, width=0, height=0):
        """
        find sizes for objects regardless of type
        :param width: (initialized with randomly large number)
        :param height: (initialized with randomly large number)
        """

        def resize(width, height, lim_width, lim_height):
            width = lim_width if width == 0 else width
            height = lim_height if height == 0 else height

            sd = {'width': lim_width, 'height': lim_height}
            d = {'width': width, 'height': height}

            if lim_width > lim_height:  # horizontal placeholder
                max_key = min(d.keys(), key=lambda k: d[k])
            else:  # vertical placeholder
                max_key = max(d.keys(), key=lambda k: d[k])

            if d[max_key] > sd[max_key]:  # check if needs resizing
                new_max_value = sd[max_key]
            else:  # doesn't need resizing
                return d['width'], d['height']

            # find counterpart
            all_keys = list(d.keys())
            all_keys.pop(all_keys.index(max_key))
            counterpart = all_keys[0]

            # adjust ratio
            d['new_{}'.format(counterpart)] = new_max_value * d[counterpart] / d[max_key]
            d['new_{}'.format(max_key)] = new_max_value

            return d['new_width'], d['new_height']

        # resize according to layout
        if layout_number == 3:  # 2 items on slide (vertical placeholders)
            if placeholder == 1:  # left item
                x = self.X
            else:  # right item
                x = self.TWO_ITEM_X
            width, height = resize(width, height, self.TWO_ITEM_WIDTH, self.ITEM_HEIGHT)

        else:  # only 1 item on slide (horizontal placeholder)
            width, height = resize(width, height, self.ONE_ITEM_WIDTH, self.ITEM_HEIGHT)
            x = max((self.SLIDE_WIDTH - width) / 2, self.X)

        return x, width, height

    def add_image(self, item, shapes, placeholder, layout_number):
        """
        add image to slide
        :param item: file name
        """
        im = Image.open(item)

        # image size
        width = im.width  # pt
        height = im.height  # pt

        x, width, height = self.object_size(placeholder, layout_number, width, height)

        # add image
        fig = BytesIO()
        im.save(fig, 'PNG')
        fig.seek(0)
        shapes.add_picture(fig, left=Pt(x), top=Pt(self.Y), width=Pt(width), height=Pt(height))
        return shapes

    def add_table(self, df, shapes, placeholder, layout_number):
        """
        add table to slide
        :param df: pandas dataframe
        """
        x, width, height = self.object_size(placeholder, layout_number)
        table = shapes.add_table(df.shape[0] + 1, df.shape[1] + 1, Pt(x), Pt(self.Y), Pt(width), Pt(height)).table

        for i, col in enumerate(df.columns):
            table.cell(0, i + 1).text = str(col)

        for i, row in enumerate(df.index):
            table.cell(i + 1, 0).text = str(row)
            table.rows[i].height = self.TABLE_ROW_HEIGHT

        for i, row in enumerate(df.index):
            table.rows[i + 1].height = self.TABLE_ROW_HEIGHT
            for e, item in enumerate(df.loc[row]):
                table.cell(i + 1, e + 1).text = str(item)
        return table

    def add_plot(self, chart_class, shapes, placeholder, layout_number):
        """
        add pptx chart to slide. currently supports line and bar plots only.
        :param chart_class: dictionary containing the dataframe and chart type.
        ex) {'df': df, 'chart_type': 'bar' or 'line'}
        """
        chart_data = ChartData()
        df = chart_class.df
        chart_data.categories = df.columns
        for i, r in enumerate(df.index):
            chart_data.add_series(str(r), df.iloc[i])

        chart_type = chart_class.chart_type
        x, width, height = self.object_size(placeholder, layout_number)

        chart = shapes.add_chart(chart_type, Pt(x), Pt(self.Y), Pt(width), Pt(height), chart_data).chart
        chart.has_legend = True
        chart.legend.include_in_layout = False
        chart.value_axis.has_major_gridlines = False
        return shapes

    def identify_content(self, content, **kwargs):
        """
        identifies content and returns appropriate functions
        """
        # bullet point
        if isinstance(content, list):
            truth_value = all([isinstance(i, str) for i in content])
            if truth_value:
                return self.add_bulletpoint(content, **kwargs)
            else:
                raise TypeError('리스트 내의 자료형이 문자열이 아닙니다')
        # image
        elif isinstance(content, str):
            try:
                Image.open(content)
                return self.add_image(content, **kwargs)
            except FileNotFoundError:
                raise FileNotFoundError('파일명을 확인해주세요')
        # table
        elif isinstance(content, pd.DataFrame):
            return self.add_table(content, **kwargs)
        else:
            try:
                # chart
                if isinstance(content.df, pd.DataFrame):
                    return self.add_plot(content, **kwargs)
                else:
                    raise TypeError('지원하지 않는 자료입니다')
            except:
                raise TypeError('지원하지 않는 자료입니다')

    def add(self, title, content1, content2=None):
        """
        main function to add content to slides.
        :param title: title of slide. ex) 'This is a title'
        :param content1: content to insert on the left
        :param content2: content to insert on the right
        """
        # slide layout
        try:
            layout_number = 3 if content2 else 1  # add layout_number 5...?
        except ValueError:
            layout_number = 1 if content2.empty else 3
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[layout_number])

        # title
        shapes = slide.shapes
        shapes.title.text = title

        # content1
        placeholder = 1
        self.identify_content(content=content1, shapes=shapes, placeholder=placeholder, layout_number=layout_number)

        # content 2
        if layout_number == 3:
            placeholder = 2
            self.identify_content(content=content2, shapes=shapes, placeholder=placeholder, layout_number=layout_number)

    def save(self, file_name):
        """
        save slides.
        :param file_name: file name. ex) 'test.pptx'
        """
        self.prs.save(file_name)
        print('"{}" 파일이 생성되었습니다.'.format(file_name))


class LineChart:
    def __init__(self, df):
        self.df = df
        self.chart_type = XL_CHART_TYPE.LINE


class BarChart:
    def __init__(self, df):
        self.df = df
        self.chart_type = XL_CHART_TYPE.COLUMN_CLUSTERED
